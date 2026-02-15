from pathlib import Path
from typing import Dict, List


PROJECT_ROOT = Path(__file__).resolve().parents[3]


def _resolve_destination(path_str: str) -> Path:
    path = Path(path_str).expanduser()
    if not path.is_absolute():
        path = (PROJECT_ROOT / path).resolve()
    else:
        path = path.resolve()
    if path.suffix.lower() != ".pptx":
        raise ValueError("pptx_writer only outputs .pptx files.")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def pptx_writer(
    file_path: str,
    slides: List[Dict[str, List[str]]],
    overwrite: bool = False,
) -> Dict[str, str]:
    if not slides:
        raise ValueError("slides must include at least one slide definition.")

    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ImportError("pptx_writer requires the 'python-pptx' package.") from exc

    destination = _resolve_destination(file_path)
    if destination.exists() and not overwrite:
        raise FileExistsError(f"File already exists (set overwrite=true to replace): {destination}")

    presentation = Presentation()

    layout = presentation.slide_layouts[1]  # Title and Content
    for slide_def in slides:
        title = slide_def.get("title", "")
        bullets = slide_def.get("bullets") or []

        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title

        body_shape = slide.shapes.placeholders[1]
        text_frame = body_shape.text_frame
        if bullets:
            text_frame.text = bullets[0]
            for bullet in bullets[1:]:
                parag = text_frame.add_paragraph()
                parag.text = bullet
                parag.level = 0
        else:
            text_frame.text = ""

    presentation.save(destination)
    return {"status": "written", "slides": str(len(slides)), "path": str(destination)}

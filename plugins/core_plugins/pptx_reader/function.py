from typing import Any, Dict, List

from services.path_utils import resolve_user_file


def pptx_reader(
    file_path: str,
    include_notes: bool = False,
    **kwargs: Any,
) -> Dict[str, Any]:
    # Backward compatibility: gracefully ignore unexpected flags (e.g., include_tables).
    _ = kwargs  # prevents unused-var lint complaints
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ImportError("pptx_reader requires the 'python-pptx' package.") from exc

    path = resolve_user_file(file_path, expected_extensions=(".pptx",))
    presentation = Presentation(path)

    slides: List[Dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                clean = shape.text.strip()
                if clean:
                    texts.append(clean)

        notes_text = ""
        if include_notes and getattr(slide, "notes_slide", None):
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()

        slides.append(
            {
                "index": index,
                "text_blocks": texts,
                "notes": notes_text if include_notes else None,
            }
        )

    return {
        "slide_count": len(slides),
        "slides": slides,
        "path": str(path),
        "notes_included": include_notes,
    }
